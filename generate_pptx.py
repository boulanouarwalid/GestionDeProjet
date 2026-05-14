from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_DARK = RGBColor(0x0B, 0x0E, 0x1A)
BLUE_NEON = RGBColor(0x00, 0xD4, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
EMERALD = RGBColor(0x00, 0xFF, 0x88)
DARK_CARD = RGBColor(0x12, 0x18, 0x2E)
GRAY = RGBColor(0x88, 0x88, 0xAA)
ORANGE = RGBColor(0xFF, 0x88, 0x00)
RED = RGBColor(0xFF, 0x33, 0x33)
YELLOW = RGBColor(0xFF, 0xCC, 0x00)
PURPLE = RGBColor(0xBB, 0x66, 0xFF)
GRID_LINE = RGBColor(0x1A, 0x22, 0x44)

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_rich_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_paragraph(tf, text, size=16, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p

def add_grid_bg(slide):
    for x in range(0, 14, 2):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(0), Inches(0.01), Inches(7.5))
        line.fill.solid()
        line.fill.fore_color.rgb = GRID_LINE
        line.line.fill.background()
    for y in range(0, 8, 2):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(y), Inches(13.333), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = GRID_LINE
        line.line.fill.background()

def add_card(slide, left, top, width, height, title, items, accent_color=BLUE_NEON):
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill_color=DARK_CARD, line_color=accent_color, line_width=1.5)
    card.shadow.inherit = False
    add_textbox(slide, left + 0.2, top + 0.15, width - 0.4, 0.4, title, font_size=14, color=accent_color, bold=True)
    y_off = top + 0.55
    for item in items:
        add_textbox(slide, left + 0.25, y_off, width - 0.5, 0.25, item, font_size=10, color=WHITE)
        y_off += 0.28

def add_badge(slide, left, top, text, color=BLUE_NEON, bg=DARK_CARD):
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, len(text)*0.12 + 0.3, 0.35, fill_color=bg, line_color=color, line_width=1)
    add_textbox(slide, left + 0.1, top + 0.02, len(text)*0.12 + 0.1, 0.3, text, font_size=9, color=color, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 1: Introduction & Contexte Métier
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1)
add_grid_bg(slide1)

add_textbox(slide1, 0.5, 0.4, 12, 0.8, "PROJECT.EVAL", font_size=44, color=BLUE_NEON, bold=True)
add_textbox(slide1, 0.5, 1.1, 12, 0.6, "Industrialisation de la Gestion de Projets", font_size=24, color=WHITE)

# Pipeline visuel - lifecycle
steps = [
    ("Init", "Initialisation", BLUE_NEON),
    ("Plan", "Planification\ndes Tâches", EMERALD),
    ("Budget", "Suivi du\nBudget", ORANGE),
    ("Report", "Rapports\nMulti-Formats", PURPLE),
]
box_w = 2.8
gap = 0.3
total_w = len(steps) * box_w + (len(steps) - 1) * gap
start_x = (13.333 - total_w) / 2
for i, (short, label, color) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    box = add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.2, box_w, 1.4, fill_color=DARK_CARD, line_color=color, line_width=2)
    add_textbox(slide1, x, 2.35, box_w, 0.5, short, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, x, 2.85, box_w, 0.7, label, font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        arrow = add_shape(slide1, MSO_SHAPE.RIGHT_ARROW, x + box_w + 0.02, 2.7, gap - 0.04, 0.4, fill_color=GRAY, line_color=None)

# Contexte cards
add_card(slide1, 0.5, 4.2, 5.8, 1.5, "CONTEXTE", [
    "Application d'entreprise pour pilotage de projets multi-départements",
    "Spécialités : DevOps, Génie Civil, Cyber, Data Science, Electricité",
], BLUE_NEON)

add_card(slide1, 7.0, 4.2, 5.8, 1.5, "PILIERS METIER", [
    "Traçabilité stricte des dépenses",
    "Alertes budgetaires en temps reel",
    "Generation de rapports multi-formats (PDF, Excel)",
], EMERALD)

# Footer
add_textbox(slide1, 0.5, 6.8, 12, 0.4, "Design Patterns en action  |  Java  Spring Boot  Angular", font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: La Problématique Technique
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2)
add_grid_bg(slide2)

add_textbox(slide2, 0.5, 0.4, 12, 0.7, "CRITICAL.ISSUES", font_size=42, color=RED, bold=True)
add_textbox(slide2, 0.5, 1.0, 12, 0.5, "Les Limites du Monolithe Refractaire", font_size=22, color=WHITE)

# Server overload icon (big red box)
server = add_shape(slide2, MSO_SHAPE.ROUNDED_RECTANGLE, 4.5, 1.8, 4.3, 2.5, fill_color=RGBColor(0x2A, 0x08, 0x08), line_color=RED, line_width=2)
add_textbox(slide2, 4.8, 2.0, 3.7, 0.5, "MONOLITHE", font_size=28, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide2, 4.8, 2.5, 3.7, 0.5, "Code Spaghetti", font_size=18, color=RGBColor(0xFF, 0x66, 0x66), alignment=PP_ALIGN.CENTER)
add_textbox(slide2, 4.8, 2.9, 3.7, 0.4, "Tout est couple partout", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# Error arrows
for angle, lbl in [(-30, "Depenses"), (30, "Budget"), (90, "Taches")]:
    rad = math.radians(angle)
    x2 = 4.5 + 2.2 + 1.8 * math.cos(rad)
    y2 = 1.8 + 1.25 + 1.8 * math.sin(rad)

# Three problem boxes
problems = [
    ("COUPLAGE FORT", "Modifier le module Depenses\ncasse le module Budget", RED),
    ("RIGIDITE", "Ajout d'une specialite = \nreescriture des controleurs", ORANGE),
    ("DETTE TECHNIQUE", "Code infrastructure (transactions,\naudits, logs) pollue le metier", YELLOW),
]
for i, (title, desc, color) in enumerate(problems):
    x = 0.8 + i * 4.2
    add_card(slide2, x, 4.8, 3.8, 1.8, title, [desc], color)

# Key quote
add_textbox(slide2, 0.5, 6.8, 12, 0.5, '"Objectif : une architecture ou le code infrastructure ne pollue jamais les regles metier"', font_size=13, color=BLUE_NEON, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3: Diagramme de Classes Global
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3)
add_grid_bg(slide3)

add_textbox(slide3, 0.5, 0.3, 12, 0.6, "ARCHITECTURE.BLUEPRINT", font_size=38, color=BLUE_NEON, bold=True)
add_textbox(slide3, 0.5, 0.8, 12, 0.4, "Cartographie Spatiale du Systeme", font_size=20, color=WHITE)

# === LAYER 1: Controller ===
add_textbox(slide3, 0.3, 1.3, 3, 0.3, "CONTROLLER", font_size=11, color=GRAY, bold=True)
ctrls = ["ProjetController", "TacheController", "BudgetController", "DepenseController"]
for i, c in enumerate(ctrls):
    x = 0.3 + i * 3.2
    add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.6, 3.0, 0.35, fill_color=DARK_CARD, line_color=BLUE_NEON, line_width=1)
    add_textbox(slide3, x, 1.63, 3.0, 0.3, c, font_size=10, color=BLUE_NEON, alignment=PP_ALIGN.CENTER)

# Arrow down
add_shape(slide3, MSO_SHAPE.DOWN_ARROW, 6.3, 1.95, 0.5, 0.3, fill_color=GRAY, line_color=None)

# === LAYER 2: Facade ===
add_textbox(slide3, 0.3, 2.3, 3, 0.3, "FACADE", font_size=11, color=EMERALD, bold=True)
facade = add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, 4.5, 2.6, 4.0, 0.4, fill_color=DARK_CARD, line_color=EMERALD, line_width=2)
add_textbox(slide3, 4.5, 2.63, 4.0, 0.35, "ProjetFacade  initialiserProjet()", font_size=11, color=EMERALD, alignment=PP_ALIGN.CENTER)

# Arrow down
add_shape(slide3, MSO_SHAPE.DOWN_ARROW, 6.3, 3.0, 0.5, 0.25, fill_color=GRAY, line_color=None)

# === LAYER 3: Services (Singleton + Proxy + Decorator + Strategy) ===
add_textbox(slide3, 0.3, 3.3, 3, 0.3, "SERVICES  @Service  @Transactional", font_size=11, color=BLUE_NEON, bold=True)
services = [
    ("ProjetService", BLUE_NEON), ("BudgetService", EMERALD),
    ("BudgetAlerteDecorator", ORANGE), ("TacheService", PURPLE),
    ("DepenseService", YELLOW),
]
for i, (srv, color) in enumerate(services):
    x = 0.3 + i * 2.55
    add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, x, 3.6, 2.35, 0.35, fill_color=DARK_CARD, line_color=color, line_width=1)
    add_textbox(slide3, x, 3.63, 2.35, 0.3, srv, font_size=9, color=color, alignment=PP_ALIGN.CENTER)

# Arrow down
add_shape(slide3, MSO_SHAPE.DOWN_ARROW, 6.3, 3.95, 0.5, 0.25, fill_color=GRAY, line_color=None)

# === LAYER 4: Strategy + Event + Template ===
add_textbox(slide3, 0.3, 4.25, 3, 0.3, "STRATEGY  /  OBSERVER  /  TEMPLATE", font_size=11, color=EMERALD, bold=True)

strat_box = add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, 0.3, 4.55, 3.8, 0.7, fill_color=DARK_CARD, line_color=EMERALD, line_width=1)
add_textbox(slide3, 0.5, 4.58, 3.6, 0.3, "Strategy: AvancementStrategy", font_size=10, color=EMERALD)
add_textbox(slide3, 0.5, 4.82, 3.6, 0.3, "  CalculSimple  |  CalculPrudent", font_size=9, color=WHITE)

obs_box = add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, 4.5, 4.55, 4.0, 0.7, fill_color=DARK_CARD, line_color=ORANGE, line_width=1)
add_textbox(slide3, 4.7, 4.58, 3.6, 0.3, "Observer: DepenseCreatedEvent", font_size=10, color=ORANGE)
add_textbox(slide3, 4.7, 4.82, 3.6, 0.3, "  DepenseService  TacheService", font_size=9, color=WHITE)

tmpl_box = add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, 8.8, 4.55, 4.0, 0.7, fill_color=DARK_CARD, line_color=PURPLE, line_width=1)
add_textbox(slide3, 9.0, 4.58, 3.6, 0.3, "Template: RapportGenerator", font_size=10, color=PURPLE)
add_textbox(slide3, 9.0, 4.82, 3.6, 0.3, "  RapportPDF  |  RapportExcel", font_size=9, color=WHITE)

# Arrow down
add_shape(slide3, MSO_SHAPE.DOWN_ARROW, 6.3, 5.25, 0.5, 0.25, fill_color=GRAY, line_color=None)

# === LAYER 5: Factory + Repository ===
add_textbox(slide3, 0.3, 5.55, 3, 0.3, "FACTORY  /  REPOSITORY", font_size=11, color=YELLOW, bold=True)

fact_box = add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, 0.3, 5.85, 6.0, 0.4, fill_color=DARK_CARD, line_color=YELLOW, line_width=1)
add_textbox(slide3, 0.5, 5.88, 5.8, 0.35, "TacheFactory (6 variantes: Cyber, DevOps, GenieCivil, Electricite, DataScience, SE)", font_size=9, color=YELLOW)

repo_box = add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, 6.8, 5.85, 6.0, 0.4, fill_color=DARK_CARD, line_color=BLUE_NEON, line_width=1)
add_textbox(slide3, 7.0, 5.88, 5.8, 0.35, "JPA Repository (Projet, Tache, Budget, Depense)", font_size=9, color=BLUE_NEON)

# === Bottom: Entity layer ===
add_shape(slide3, MSO_SHAPE.DOWN_ARROW, 6.3, 6.25, 0.5, 0.25, fill_color=GRAY, line_color=None)
add_textbox(slide3, 0.3, 6.55, 3, 0.3, "ENTITIES  (JPA)", font_size=11, color=WHITE, bold=True)
entities = ["Projet  (Cloneable)", "Tache", "Budget", "Depense"]
for i, e in enumerate(entities):
    x = 0.8 + i * 3.2
    add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, x, 6.85, 2.8, 0.35, fill_color=DARK_CARD, line_color=WHITE, line_width=1)
    add_textbox(slide3, x, 6.88, 2.8, 0.3, e, font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)

# Legend
add_textbox(slide3, 10.5, 1.3, 2.5, 0.3, "LEGENDE", font_size=10, color=GRAY)
legend_items = [("Services", BLUE_NEON), ("Factory", YELLOW), ("Facade", EMERALD), ("Decorator", ORANGE), ("Observer", ORANGE), ("Template", PURPLE)]
for i, (lbl, c) in enumerate(legend_items):
    add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, 10.5, 1.6 + i * 0.3, 0.3, 0.2, fill_color=c, line_color=None)
    add_textbox(slide3, 10.9, 1.58, 1.5, 0.25, lbl, font_size=8, color=WHITE)

# ============================================================
# SLIDE 4: Patterns Part 1 - Création & Structure
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4)
add_grid_bg(slide4)

add_textbox(slide4, 0.5, 0.3, 12, 0.6, "01 // ARCHITECTURE.PATTERNS", font_size=36, color=BLUE_NEON, bold=True)
add_textbox(slide4, 0.5, 0.8, 12, 0.4, "Instanciation & Isolation", font_size=20, color=WHITE)

# --- CREATION section ---
add_textbox(slide4, 0.3, 1.2, 3, 0.4, "--- CREATION ---", font_size=16, color=EMERALD, bold=True)

# Singleton
add_card(slide4, 0.3, 1.7, 6.0, 1.4, "SINGLETON", [
    "8 classes @Service : 1 seule instance par conteneur Spring IoC",
    "ProjetService  TacheService  BudgetService  DepenseService",
    "BudgetAlerteDecorator  ProjetFacade  RapportPDF  RapportExcel",
], BLUE_NEON)
add_badge(slide4, 5.0, 1.7, "Spring default", BLUE_NEON)

# Factory Method
add_card(slide4, 6.6, 1.7, 6.2, 1.4, "FACTORY METHOD", [
    "TacheFactory interface : creerTache(nom, projet)",
    "6 variantes : Cyber, DevOps, GenieCivil, Electricite, DataScience, SE",
    "Evite le mot-cle new  —  Dispatch dynamique par type",
], YELLOW)
add_badge(slide4, 11.3, 1.7, "6 factories", YELLOW)

# Builder
add_card(slide4, 0.3, 3.3, 6.0, 1.2, "BUILDER", [
    "ProjetDTO  TacheDTO  DepenseDTO : @Builder (Lombok)",
    "Construction fluide et immuable : ProjetDTO.builder().nomProjet(...).build()",
    "Validation integree via @NotBlank  @Size  @NotNull",
], EMERALD)
add_badge(slide4, 5.0, 3.3, "Lombok", EMERALD)

# Prototype
add_card(slide4, 6.6, 3.3, 6.2, 1.2, "PROTOTYPE", [
    "Projet implements Cloneable + clone()",
    "reset Id=null, Taches=null, Budgets=null pour copie propre",
    "Endpoint POST /{id}/cloner  —  Duplication en RAM",
], BLUE_NEON)
add_badge(slide4, 11.3, 3.3, "Cloneable", BLUE_NEON)

# --- STRUCTURE section ---
add_textbox(slide4, 0.3, 4.7, 3, 0.4, "--- STRUCTURE ---", font_size=16, color=ORANGE, bold=True)

# Facade
add_card(slide4, 0.3, 5.2, 6.0, 1.0, "FACADE", [
    "ProjetFacade unifie 4 services en 1 appel",
    "initialiserProjet(ProjetDTO, Map<TacheDTO, DepenseDTO>, BudgetDTO)",
    "Masque la complexite de creation projet + budget + taches + depenses",
], EMERALD)
add_badge(slide4, 5.0, 5.2, "1 methode", EMERALD)

# Proxy
add_card(slide4, 6.6, 5.2, 6.2, 1.0, "PROXY", [
    "@Transactional sur 5 services  —  Proxy Spring AOP",
    "beginTransaction / commit / rollback automatique",
    "Code metier zero infrastructure transactionnelle",
], BLUE_NEON)
add_badge(slide4, 11.3, 5.2, "@Transactional", BLUE_NEON)

# Decorator
add_card(slide4, 0.3, 6.4, 12.5, 0.9, "DECORATOR", [
    "BudgetAlerteDecorator implements IBudgetService",
    "Wrapper around BudgetService : ajoute alerte si budget > 500 000 MAD",
    "Audit des transactions a la volee sans modifier BudgetService original",
], ORANGE)
add_badge(slide4, 11.3, 6.4, "alerte", ORANGE)

# ============================================================
# SLIDE 5: Patterns Part 2 - Comportement & Réactivité
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5)
add_grid_bg(slide5)

add_textbox(slide5, 0.5, 0.3, 12, 0.6, "02 // ARCHITECTURE.PATTERNS", font_size=36, color=PURPLE, bold=True)
add_textbox(slide5, 0.5, 0.8, 12, 0.4, "Algorithmes & Flux Reactifs", font_size=20, color=WHITE)

# Strategy
add_card(slide5, 0.3, 1.5, 4.0, 2.2, "STRATEGY", [
    "AvancementStrategy interface",
    "  calculer(prevu, consomme)",
    "",
    "CalculSimple : (consomme / prevu) * 100",
    "CalculPrudent : (consomme * 1.1 / prevu) * 100",
    "",
    "Injection via @Qualifier(\"calculSimple\")",
    "Interchangeable sans modifier BudgetService",
], EMERALD)
add_badge(slide5, 3.0, 1.5, "@Qualifier", EMERALD)

# Observer - schema
obs_card = add_card(slide5, 4.6, 1.5, 4.0, 2.2, "OBSERVER", [
    "Event : DepenseCreatedEvent",
    "",
    "Publisher : DepenseService",
    "  eventPublisher.publishEvent(event)",
    "",
    "Listener : TacheService",
    "  @EventListener  update Budget",
    "",
    "Couplage = ZERO  (canal neutre)",
], ORANGE)
add_badge(slide5, 7.3, 1.5, "event-driven", ORANGE)

# Observer flow diagram
flow_y = 2.6
add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, 4.8, flow_y, 1.3, 0.4, fill_color=DARK_CARD, line_color=YELLOW, line_width=1)
add_textbox(slide5, 4.8, flow_y + 0.05, 1.3, 0.3, "DepenseService", font_size=8, color=YELLOW, alignment=PP_ALIGN.CENTER)
add_shape(slide5, MSO_SHAPE.RIGHT_ARROW, 6.15, flow_y + 0.05, 0.5, 0.3, fill_color=ORANGE, line_color=None)
add_textbox(slide5, 6.55, flow_y + 0.05, 0.6, 0.3, "Event", font_size=8, color=ORANGE, alignment=PP_ALIGN.CENTER)
add_shape(slide5, MSO_SHAPE.RIGHT_ARROW, 7.2, flow_y + 0.05, 0.5, 0.3, fill_color=ORANGE, line_color=None)
add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, 7.75, flow_y, 1.3, 0.4, fill_color=DARK_CARD, line_color=PURPLE, line_width=1)
add_textbox(slide5, 7.75, flow_y + 0.05, 1.3, 0.3, "TacheService", font_size=8, color=PURPLE, alignment=PP_ALIGN.CENTER)
add_shape(slide5, MSO_SHAPE.DOWN_ARROW, 8.3, flow_y + 0.35, 0.3, 0.25, fill_color=WHITE, line_color=None)
add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, 8.0, flow_y + 0.65, 0.9, 0.3, fill_color=DARK_CARD, line_color=WHITE, line_width=1)
add_textbox(slide5, 8.0, flow_y + 0.68, 0.9, 0.25, "Budget", font_size=7, color=WHITE, alignment=PP_ALIGN.CENTER)

# Template Method
add_card(slide5, 8.9, 1.5, 4.0, 2.2, "TEMPLATE METHOD", [
    "RapportGenerator (abstract)",
    "",
    "Template : genererRapportComplet() [final]",
    "  1. extraireDonnees()  [concret]",
    "  2. formaterFichier()  [abstract]",
    "  3. Envoi au client",
    "",
    "RapportPDFService  |  RapportExcelService",
], PURPLE)
add_badge(slide5, 11.6, 1.5, "abstract", PURPLE)

# Bottom: detailed schema for Strategy + Template
add_textbox(slide5, 0.3, 4.2, 12, 0.4, "SCHEMA D'EXECUTION", font_size=14, color=BLUE_NEON, bold=True)

# Flow: Client -> Controller -> Facade -> Services -> ...
add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5, 4.8, 1.2, 0.5, fill_color=DARK_CARD, line_color=WHITE, line_width=1)
add_textbox(slide5, 0.5, 4.85, 1.2, 0.4, "Client", font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)
add_shape(slide5, MSO_SHAPE.RIGHT_ARROW, 1.75, 4.95, 0.4, 0.2, fill_color=BLUE_NEON, line_color=None)
add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, 2.2, 4.8, 1.5, 0.5, fill_color=DARK_CARD, line_color=BLUE_NEON, line_width=1)
add_textbox(slide5, 2.2, 4.85, 1.5, 0.4, "Controller", font_size=11, color=BLUE_NEON, alignment=PP_ALIGN.CENTER)
add_shape(slide5, MSO_SHAPE.RIGHT_ARROW, 3.75, 4.95, 0.4, 0.2, fill_color=EMERALD, line_color=None)
add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, 4.2, 4.8, 1.5, 0.5, fill_color=DARK_CARD, line_color=EMERALD, line_width=1)
add_textbox(slide5, 4.2, 4.85, 1.5, 0.4, "ProjetFacade", font_size=11, color=EMERALD, alignment=PP_ALIGN.CENTER)

# Branching to services
add_shape(slide5, MSO_SHAPE.RIGHT_ARROW, 5.75, 4.95, 0.4, 0.2, fill_color=GRAY, line_color=None)
add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, 6.2, 4.8, 1.5, 0.5, fill_color=DARK_CARD, line_color=ORANGE, line_width=1)
add_textbox(slide5, 6.2, 4.85, 1.5, 0.4, "Services", font_size=11, color=ORANGE, alignment=PP_ALIGN.CENTER)

# Observer branch
add_shape(slide5, MSO_SHAPE.DOWN_ARROW, 6.0, 5.35, 0.3, 0.3, fill_color=ORANGE, line_color=None)
add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, 5.5, 5.7, 1.5, 0.5, fill_color=DARK_CARD, line_color=ORANGE, line_width=1)
add_textbox(slide5, 5.5, 5.75, 1.5, 0.4, "Event Bus", font_size=11, color=ORANGE, alignment=PP_ALIGN.CENTER)
add_shape(slide5, MSO_SHAPE.RIGHT_ARROW, 7.05, 5.85, 0.4, 0.2, fill_color=PURPLE, line_color=None)
add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, 7.5, 5.7, 1.5, 0.5, fill_color=DARK_CARD, line_color=PURPLE, line_width=1)
add_textbox(slide5, 7.5, 5.75, 1.5, 0.4, "Listener", font_size=11, color=PURPLE, alignment=PP_ALIGN.CENTER)

# Strategy branch
add_shape(slide5, MSO_SHAPE.UP_ARROW, 8.6, 5.35, 0.3, 0.3, fill_color=EMERALD, line_color=None)
add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, 8.0, 4.3, 2.0, 0.5, fill_color=DARK_CARD, line_color=EMERALD, line_width=1)
add_textbox(slide5, 8.0, 4.35, 2.0, 0.4, "CalculSimple", font_size=9, color=EMERALD, alignment=PP_ALIGN.CENTER)
add_textbox(slide5, 8.0, 4.6, 2.0, 0.2, "ou CalculPrudent", font_size=7, color=GRAY, alignment=PP_ALIGN.CENTER)

# Template branch
add_shape(slide5, MSO_SHAPE.DOWN_ARROW, 10.0, 5.35, 0.3, 0.3, fill_color=PURPLE, line_color=None)
add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, 9.3, 5.7, 2.0, 0.5, fill_color=DARK_CARD, line_color=PURPLE, line_width=1)
add_textbox(slide5, 9.3, 5.75, 2.0, 0.4, "RapportPDF", font_size=9, color=PURPLE, alignment=PP_ALIGN.CENTER)
add_textbox(slide5, 9.3, 6.0, 2.0, 0.2, "ou RapportExcel", font_size=7, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: Démonstration de l'Application
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6)
add_grid_bg(slide6)

add_textbox(slide6, 0.5, 0.3, 12, 0.6, "RUNTIME.LIVE_DEMO", font_size=38, color=EMERALD, bold=True)
add_textbox(slide6, 0.5, 0.8, 12, 0.4, "La Preuve par le Code", font_size=20, color=WHITE)

# Terminal window mockup
term_bg = RGBColor(0x0D, 0x11, 0x1D)
term = add_shape(slide6, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5, 1.5, 12.3, 4.5, fill_color=term_bg, line_color=BLUE_NEON, line_width=1.5)

# Terminal title bar
add_shape(slide6, MSO_SHAPE.RECTANGLE, 0.5, 1.5, 12.3, 0.35, fill_color=RGBColor(0x1A, 0x22, 0x44), line_color=None)
add_textbox(slide6, 0.8, 1.52, 5, 0.3, "PROJECT.EVAL - Test Execution", font_size=11, color=GRAY, bold=True)

# Terminal content - simulated log output
log_lines = [
    ("$  mvn test -Dtest=GestionProjetApplicationTests", EMERALD),
    ("", None),
    ("[INFO]  T E S T S", BLUE_NEON),
    ("[INFO] Running com.controleproject.GestionProjetApplicationTests", GRAY),
    ("", None),
    ("  [PROXY] @Transactional: begin transaction", BLUE_NEON),
    ("  [FACADE] ProjetFacade.initialiserProjet() orchestre 4 services...", EMERALD),
    ("  [FACTORY] TacheGenieCivilFactory.creerTache() -> Tache type=GENIE_CIVIL", YELLOW),
    ("  [OBSERVER] DepenseCreatedEvent publie -> TacheService recoit", ORANGE),
    ("  [OBSERVER] Budget.montantConsomme mis a jour: 15000.0", ORANGE),
    ("  [DECORATOR] ALERTE : Budget depasse 500 000 MAD !", RED),
    ("  [TEMPLATE] RapportPDFService.genererRapportComplet() -> format PDF", PURPLE),
    ("  [STRATEGY] CalculSimple: taux avancement = 45.3%", EMERALD),
    ("", None),
    ("[INFO] Tests run: 10, Failures: 0, Errors: 0, Skipped: 0", EMERALD),
    ("[INFO] BUILD SUCCESS", EMERALD),
]

y = 2.0
for text, color in log_lines:
    if text == "":
        y += 0.12
        continue
    c = color if color else WHITE
    add_textbox(slide6, 0.8, y, 11.5, 0.25, text, font_size=9, color=c, font_name="Consolas")
    y += 0.28

# Bottom: Success badge + instructions
success = add_shape(slide6, MSO_SHAPE.ROUNDED_RECTANGLE, 3.5, 6.2, 6.3, 0.7, fill_color=RGBColor(0x08, 0x2A, 0x10), line_color=EMERALD, line_width=2)
add_textbox(slide6, 3.5, 6.25, 6.3, 0.6, "TESTS PASSES  |  BUILD SUCCESS  |  ZONE VERITE", font_size=18, color=EMERALD, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide6, 0.5, 7.0, 12, 0.3, "Les 10 patterns s'executent a la chaine  |  Barre verte JUnit = validation architecturale", font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# Save
output_path = "C:\\Users\\Waaleed Blr Anir\\OneDrive\\Bureau\\projects\\GestionDeProjet\\GestionDeProjet_Patterns.pptx"
prs.save(output_path)
print(f"PPTX created: {output_path}")
